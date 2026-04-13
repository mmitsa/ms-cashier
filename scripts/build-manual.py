#!/usr/bin/env python3
"""Build Cashier User Manual PPTX — A4 Landscape with screenshots."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# A4 Landscape dimensions
A4_W = Cm(29.7)
A4_H = Cm(21.0)

# Brand colors
BRAND_PRIMARY = RGBColor(0x4F, 0x46, 0xE5)  # Indigo
BRAND_DARK = RGBColor(0x1E, 0x1B, 0x4B)
BRAND_LIGHT = RGBColor(0xE0, 0xE7, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x1F, 0x29, 0x37)

SCREENSHOTS_DIR = os.path.join(os.path.dirname(__file__), '..', 'screenshots')

def add_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=DARK, bold=False, alignment=PP_ALIGN.RIGHT, font_name='Arial'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_step_number(slide, left, top, number):
    """Add a circular step number badge."""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Cm(1.8), Cm(1.8)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = BRAND_PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

def add_screenshot(slide, img_path, left, top, max_width, max_height=None):
    """Add screenshot image scaled to fit."""
    if not os.path.exists(img_path):
        print(f"  WARNING: {img_path} not found")
        return
    from PIL import Image
    img = Image.open(img_path)
    w, h = img.size
    ratio = w / h

    width = max_width
    height = int(width / ratio)
    if max_height and height > max_height:
        height = max_height
        width = int(height * ratio)

    slide.shapes.add_picture(img_path, left, top, width, height)

def build():
    prs = Presentation()
    prs.slide_width = A4_W
    prs.slide_height = A4_H

    slides_data = [
        # (screenshot, step_num, title, subtitle, bullets)
        {
            'type': 'cover',
            'title': 'دليل مستخدم الكاشير',
            'subtitle': 'MS Cashier — نظام نقاط البيع المتكامل',
            'footer': 'دليل تشغيلي خطوة بخطوة لموظف الكاشير'
        },
        {
            'type': 'step',
            'img': '08-login.png',
            'step': 1,
            'title': 'تسجيل الدخول',
            'bullets': [
                'افتح المتصفح وادخل على رابط النظام',
                'أدخل اسم المستخدم الخاص بك',
                'أدخل كلمة المرور',
                'اضغط زر "تسجيل الدخول"',
                'سيتم توجيهك مباشرة لشاشة نقطة البيع',
            ]
        },
        {
            'type': 'step',
            'img': 'check-state.png',
            'step': 2,
            'title': 'شاشة نقطة البيع الرئيسية',
            'bullets': [
                'الجانب الأيمن: قائمة المنتجات مع البحث والتصنيفات',
                'الجانب الأيسر: سلة المشتريات (الفاتورة)',
                'الشريط العلوي: اسم الكاشير والساعة',
                'يمكنك البحث عن المنتج بالاسم أو مسح الباركود',
            ]
        },
        {
            'type': 'step',
            'img': '06-pos-with-cart.png',
            'step': 3,
            'title': 'إضافة المنتجات للفاتورة',
            'bullets': [
                'اضغط على أي منتج لإضافته للسلة',
                'اضغط مرة أخرى لزيادة الكمية',
                'استخدم أزرار + و - لتعديل الكمية',
                'اضغط على سلة المهملات لحذف المنتج',
                'الإجمالي يتحدث تلقائياً في أسفل السلة',
            ]
        },
        {
            'type': 'step',
            'img': '07-category-filter.png',
            'step': 4,
            'title': 'تصفية المنتجات حسب التصنيف',
            'bullets': [
                'اضغط على اسم التصنيف لعرض منتجاته فقط',
                'مثال: "مشروبات" تعرض المشروبات فقط',
                'اضغط "الكل" للعودة لجميع المنتجات',
                'يظهر عدد المنتجات بجانب كل تصنيف',
            ]
        },
        {
            'type': 'step',
            'img': '08-customer-picker.png',
            'step': 5,
            'title': 'اختيار العميل (اختياري)',
            'bullets': [
                'اضغط على "عميل نقدي" لتغيير العميل',
                'ابحث عن العميل بالاسم أو رقم الهاتف',
                'اختر العميل من القائمة',
                'مطلوب عند البيع الآجل (على الحساب)',
                'اتركه "عميل نقدي" للبيع النقدي العادي',
            ]
        },
        {
            'type': 'step',
            'img': '09-cash-payment.png',
            'step': 6,
            'title': 'إتمام عملية البيع — الدفع',
            'bullets': [
                'اختر طريقة الدفع: كاش / مدى / تحويل / آجل',
                'عند اختيار "كاش" تظهر نافذة الدفع',
                'أدخل المبلغ المدفوع من العميل',
                'يحسب النظام الباقي تلقائياً',
                'اضغط "تأكيد البيع" لإتمام العملية',
            ]
        },
        {
            'type': 'step',
            'img': '06-pos-with-cart.png',
            'step': 7,
            'title': 'طرق الدفع المتاحة',
            'bullets': [
                'كاش (أخضر): الدفع النقدي — الأكثر استخداماً',
                'مدى / بطاقة (أزرق): الدفع بالبطاقة عبر الجهاز',
                'تحويل (بنفسجي): التحويل البنكي',
                'آجل (برتقالي): على الحساب — يتطلب اختيار عميل',
                'بعد التأكيد: تظهر نافذة النجاح مع رقم الفاتورة',
            ]
        },
        {
            'type': 'step',
            'img': '10-sale-success.png',
            'step': 8,
            'title': 'بعد إتمام البيع',
            'bullets': [
                'تظهر رسالة نجاح مع رقم الفاتورة',
                'اضغط "طباعة" لطباعة الإيصال',
                'اضغط "فاتورة جديدة" لبدء عملية بيع جديدة',
                'يتم خصم الكميات من المخزون تلقائياً',
                'يمكن مراجعة الفواتير من قائمة "المبيعات"',
            ]
        },
        {
            'type': 'tips',
            'title': 'نصائح وملاحظات مهمة',
            'sections': [
                ('اختصارات سريعة', [
                    'مسح الباركود بالقارئ يضيف المنتج مباشرة',
                    'اضغط "ملء الشاشة" لتوسيع واجهة نقطة البيع',
                    'اضغط "خصم" لإضافة خصم على الفاتورة',
                    'اضغط "مسح" لحذف جميع المنتجات من السلة',
                ]),
                ('ملاحظات مهمة', [
                    'تأكد من اتصال الإنترنت قبل بدء العمل',
                    'راجع الإجمالي قبل تأكيد البيع',
                    'لا تشارك بيانات تسجيل الدخول مع أحد',
                    'عند وجود مشكلة تواصل مع مدير النظام',
                ]),
            ]
        },
    ]

    for i, data in enumerate(slides_data):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        if data['type'] == 'cover':
            # === COVER SLIDE ===
            add_bg(slide, BRAND_DARK)

            # Logo area
            logo_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Cm(12.5), Cm(3.5), Cm(4.7), Cm(2.2)
            )
            logo_shape.fill.solid()
            logo_shape.fill.fore_color.rgb = BRAND_PRIMARY
            logo_shape.line.fill.background()
            tf = logo_shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = 'MS'
            p.font.size = Pt(36)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            add_textbox(slide, Cm(2), Cm(7), Cm(25.7), Cm(2.5),
                       data['title'], font_size=44, color=WHITE, bold=True,
                       alignment=PP_ALIGN.CENTER)
            add_textbox(slide, Cm(2), Cm(10), Cm(25.7), Cm(1.5),
                       data['subtitle'], font_size=22, color=BRAND_LIGHT,
                       alignment=PP_ALIGN.CENTER)
            add_textbox(slide, Cm(2), Cm(13), Cm(25.7), Cm(1.2),
                       data['footer'], font_size=16, color=GRAY,
                       alignment=PP_ALIGN.CENTER)

            # Version / date
            add_textbox(slide, Cm(2), Cm(17), Cm(25.7), Cm(1),
                       'الإصدار 1.0 — أبريل 2026', font_size=12, color=GRAY,
                       alignment=PP_ALIGN.CENTER)

        elif data['type'] == 'step':
            # === STEP SLIDE ===
            add_bg(slide, WHITE)

            # Top bar
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), A4_W, Cm(1.8)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = BRAND_PRIMARY
            top_bar.line.fill.background()

            add_textbox(slide, Cm(1), Cm(0.3), Cm(27.7), Cm(1.2),
                       f'الخطوة {data["step"]}', font_size=14, color=WHITE,
                       bold=True, alignment=PP_ALIGN.LEFT)

            add_textbox(slide, Cm(1), Cm(0.3), Cm(27.7), Cm(1.2),
                       'دليل مستخدم الكاشير — MS Cashier', font_size=11,
                       color=BRAND_LIGHT, alignment=PP_ALIGN.RIGHT)

            # Step number badge
            add_step_number(slide, Cm(25.5), Cm(2.5), data['step'])

            # Title
            add_textbox(slide, Cm(15), Cm(2.5), Cm(10), Cm(1.5),
                       data['title'], font_size=28, color=BRAND_DARK,
                       bold=True, alignment=PP_ALIGN.RIGHT)

            # Bullets
            txBox = slide.shapes.add_textbox(Cm(15), Cm(4.5), Cm(12.5), Cm(14))
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(data['bullets']):
                if j > 0:
                    p = tf.add_paragraph()
                else:
                    p = tf.paragraphs[0]
                p.text = f'  {bullet}'
                p.font.size = Pt(14)
                p.font.color.rgb = DARK
                p.font.name = 'Arial'
                p.alignment = PP_ALIGN.RIGHT
                p.space_after = Pt(10)
                # Add bullet marker
                p.text = f'◂  {bullet}'

            # Screenshot
            img_path = os.path.join(SCREENSHOTS_DIR, data['img'])
            if os.path.exists(img_path):
                add_screenshot(slide, img_path, Cm(1), Cm(2.8), Cm(13.5), Cm(16))
            else:
                print(f"  Missing: {data['img']}")

        elif data['type'] == 'tips':
            # === TIPS SLIDE ===
            add_bg(slide, WHITE)

            # Top bar
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), A4_W, Cm(1.8)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = BRAND_PRIMARY
            top_bar.line.fill.background()

            add_textbox(slide, Cm(1), Cm(0.3), Cm(27.7), Cm(1.2),
                       'نصائح وملاحظات', font_size=14, color=WHITE,
                       bold=True, alignment=PP_ALIGN.LEFT)
            add_textbox(slide, Cm(1), Cm(0.3), Cm(27.7), Cm(1.2),
                       'دليل مستخدم الكاشير — MS Cashier', font_size=11,
                       color=BRAND_LIGHT, alignment=PP_ALIGN.RIGHT)

            # Title
            add_textbox(slide, Cm(2), Cm(2.5), Cm(25.7), Cm(1.5),
                       data['title'], font_size=32, color=BRAND_DARK,
                       bold=True, alignment=PP_ALIGN.CENTER)

            # Two columns of tips
            col_x = [Cm(15.5), Cm(1.5)]
            for col_idx, (section_title, items) in enumerate(data['sections']):
                x = col_x[col_idx]

                # Section card background
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(4.8), Cm(12.8), Cm(13)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
                card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

                add_textbox(slide, x + Cm(0.5), Cm(5.3), Cm(11.8), Cm(1.2),
                           section_title, font_size=20, color=BRAND_PRIMARY,
                           bold=True, alignment=PP_ALIGN.RIGHT)

                txBox = slide.shapes.add_textbox(x + Cm(0.5), Cm(7), Cm(11.8), Cm(10))
                tf = txBox.text_frame
                tf.word_wrap = True
                for j, item in enumerate(items):
                    if j > 0:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    p.text = f'✓  {item}'
                    p.font.size = Pt(14)
                    p.font.color.rgb = DARK
                    p.font.name = 'Arial'
                    p.alignment = PP_ALIGN.RIGHT
                    p.space_after = Pt(14)

        print(f"  Slide {i+1}: {data.get('title', data['type'])}")

    output_path = os.path.join(SCREENSHOTS_DIR, '..', 'دليل-مستخدم-الكاشير.pptx')
    prs.save(output_path)
    print(f"\n✅ Saved: {os.path.abspath(output_path)}")
    return output_path

if __name__ == '__main__':
    build()
